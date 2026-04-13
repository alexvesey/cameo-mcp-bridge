"""Dry-run rubric workflow helpers for assignment/package validation.

This module stays intentionally light on bridge mutation. It reuses the
methodology registry/service layer to build rubric expectations and then
projects those expectations into reviewable plan objects:

- validate assignment/package
- export required diagrams
- assemble PPT/PDF
- compare against an expected artifact list
"""

from __future__ import annotations

import base64
from dataclasses import asdict
from pathlib import Path
import re
from typing import Any, Mapping, Sequence

from PIL import Image

from cameo_mcp import client as default_bridge_client
from cameo_mcp.methodology import registry, service as methodology_service


def _to_plain(value: Any) -> Any:
    if hasattr(value, "__dataclass_fields__"):
        return {key: _to_plain(item) for key, item in asdict(value).items()}
    if isinstance(value, dict):
        return {key: _to_plain(item) for key, item in value.items()}
    if isinstance(value, (list, tuple)):
        return [_to_plain(item) for item in value]
    return value


def _optional_str(value: Any) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    return text or None


def _normalize_text(value: Any) -> str:
    return " ".join(str(value or "").split())


def _normalize_kind(value: Any) -> str:
    return re.sub(r"[\s_-]+", "", _normalize_text(value)).casefold()


def _slugify(value: Any) -> str:
    text = _normalize_text(value).casefold()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text or "artifact"


def _is_diagram_kind(kind: Any) -> bool:
    normalized = _normalize_kind(kind)
    return normalized in {"bdd", "ibd"} or "diagram" in normalized


def _artifact_key(payload: Mapping[str, Any]) -> str:
    for key in ("key", "artifactKey", "element_id", "elementId", "id"):
        value = payload.get(key)
        if value is not None and str(value).strip():
            return str(value).strip()
    raise ValueError(f"Artifact is missing a stable key: {payload!r}")


def _coerce_artifact(payload: Mapping[str, Any] | Any, *, role: str) -> dict[str, Any]:
    if hasattr(payload, "__dataclass_fields__"):
        payload = asdict(payload)
    if not isinstance(payload, Mapping):
        raise TypeError(f"{role} artifact must be a mapping or dataclass, got {type(payload)!r}")

    properties = payload.get("properties") or {}
    if not isinstance(properties, Mapping):
        properties = {}

    stereotypes = payload.get("stereotypes") or ()
    if isinstance(stereotypes, str):
        stereotypes = (stereotypes,)

    return {
        "key": _artifact_key(payload),
        "kind": _normalize_text(payload.get("kind") or payload.get("type") or payload.get("humanType") or ""),
        "name": _optional_str(payload.get("name") or payload.get("elementName")),
        "parent_key": _optional_str(payload.get("parent_key") or payload.get("parentKey") or payload.get("ownerId")),
        "element_id": _optional_str(payload.get("element_id") or payload.get("elementId")),
        "stereotypes": tuple(_normalize_text(item) for item in stereotypes if _normalize_text(item)),
        "properties": dict(properties),
        "recipe_id": _optional_str(payload.get("recipe_id") or payload.get("recipeId")),
        "source_recipe_id": _optional_str(payload.get("source_recipe_id") or payload.get("sourceRecipeId")),
    }


def _build_expected_artifacts(
    pack_id: str,
    *,
    recipe_id: str | None = None,
) -> list[dict[str, Any]]:
    pack = registry.get_pack(pack_id)
    recipes = (pack.recipe(recipe_id),) if recipe_id is not None else pack.recipes

    expected: list[dict[str, Any]] = []
    seen_keys: set[str] = set()
    for recipe in recipes:
        for requirement in recipe.required_artifacts:
            if requirement.key in seen_keys:
                continue
            expected.append(
                {
                    "key": requirement.key,
                    "kind": requirement.kind,
                    "name": requirement.name,
                    "parent_key": requirement.parent_key,
                    "stereotypes": list(requirement.stereotypes),
                    "properties": dict(requirement.properties),
                    "recipe_id": recipe.id,
                    "recipe_name": recipe.title,
                    "phase_id": recipe.phase_id,
                }
            )
            seen_keys.add(requirement.key)
    return expected


def _expected_diagrams(expected_artifacts: Sequence[Mapping[str, Any]]) -> list[dict[str, Any]]:
    return [
        dict(item)
        for item in expected_artifacts
        if _is_diagram_kind(item.get("kind"))
    ]


def _comparison_status(reasons: Sequence[str]) -> str:
    if not reasons:
        return "match"
    if any("missing" in reason for reason in reasons):
        return "missing"
    if any("kind" in reason for reason in reasons):
        return "kind_mismatch"
    if any("name" in reason for reason in reasons):
        return "name_mismatch"
    if any("parent" in reason for reason in reasons):
        return "parent_mismatch"
    if any("stereotype" in reason for reason in reasons):
        return "stereotype_mismatch"
    if any("property" in reason for reason in reasons):
        return "property_mismatch"
    return "mismatch"


def _suggested_actions(expected: Mapping[str, Any], actual: Mapping[str, Any] | None, reasons: Sequence[str]) -> list[dict[str, Any]]:
    actions: list[dict[str, Any]] = []
    if actual is None:
        actions.append(
            {
                "action": "create_artifact",
                "targetKey": expected["key"],
                "targetKind": expected["kind"],
                "preview": (
                    f"Create {expected['kind']} '{expected.get('name') or expected['key']}'"
                    + (
                        f" under '{expected['parent_key']}'"
                        if expected.get("parent_key")
                        else ""
                    )
                    + "."
                ),
            }
        )
        return actions

    if any("kind" in reason for reason in reasons):
        actions.append(
            {
                "action": "replace_or_retype_artifact",
                "targetKey": expected["key"],
                "targetKind": expected["kind"],
                "preview": (
                    f"Replace or retype '{actual.get('name') or actual['key']}' "
                    f"to match {expected['kind']}."
                ),
            }
        )
    if any("name" in reason for reason in reasons):
        actions.append(
            {
                "action": "rename_artifact",
                "targetKey": expected["key"],
                "preview": f"Rename '{actual.get('name') or actual['key']}' to '{expected.get('name') or expected['key']}'.",
            }
        )
    if any("parent" in reason for reason in reasons):
        actions.append(
            {
                "action": "reparent_artifact",
                "targetKey": expected["key"],
                "preview": (
                    f"Move '{actual.get('name') or actual['key']}' under "
                    f"'{expected.get('parent_key') or 'the expected package'}'."
                ),
            }
        )
    if any("stereotype" in reason for reason in reasons):
        actions.append(
            {
                "action": "update_stereotypes",
                "targetKey": expected["key"],
                "preview": f"Update stereotypes on '{actual.get('name') or actual['key']}' to match the rubric.",
            }
        )
    if any("property" in reason for reason in reasons):
        actions.append(
            {
                "action": "update_properties",
                "targetKey": expected["key"],
                "preview": f"Update properties on '{actual.get('name') or actual['key']}' to match the rubric.",
            }
        )
    return actions


def compare_against_expected_artifact_list(
    expected_artifacts: Sequence[Mapping[str, Any] | Any],
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
) -> dict[str, Any]:
    """Compare current artifacts against an expected rubric list.

    The function is intentionally dry-run oriented. It returns a stable diff and
    previewable patch plan instead of mutating the model.
    """
    expected_items = [_coerce_artifact(item, role="expected") for item in expected_artifacts]
    actual_items = [_coerce_artifact(item, role="current") for item in (current_artifacts or ())]

    expected_index = {item["key"]: item for item in expected_items}
    actual_index = {item["key"]: item for item in actual_items}

    entries: list[dict[str, Any]] = []
    patch_plan: list[dict[str, Any]] = []

    for expected in expected_items:
        actual = actual_index.get(expected["key"])
        if actual is None:
            reasons = ["missing artifact"]
            entry = {
                "key": expected["key"],
                "status": "missing",
                "reasons": reasons,
                "expected": expected,
                "actual": None,
                "fixable": True,
                "suggestedActions": _suggested_actions(expected, None, reasons),
            }
            entries.append(entry)
            patch_plan.extend(
                {
                    "stepId": f"rubric:{len(patch_plan) + index}",
                    **action,
                    "status": "preview",
                    "reason": "missing artifact",
                }
                for index, action in enumerate(entry["suggestedActions"], start=1)
            )
            continue

        reasons: list[str] = []
        if _normalize_kind(actual["kind"]) != _normalize_kind(expected["kind"]):
            reasons.append("kind mismatch")
        if expected.get("name") is not None and _normalize_text(actual.get("name")) != _normalize_text(expected.get("name")):
            reasons.append("name mismatch")
        if expected.get("parent_key") is not None and _optional_str(actual.get("parent_key")) != _optional_str(expected.get("parent_key")):
            reasons.append("parent mismatch")
        expected_stereotypes = {str(item).casefold() for item in expected.get("stereotypes", ()) if str(item).strip()}
        actual_stereotypes = {str(item).casefold() for item in actual.get("stereotypes", ()) if str(item).strip()}
        if expected_stereotypes and not expected_stereotypes.issubset(actual_stereotypes):
            reasons.append("stereotype mismatch")
        expected_properties = dict(expected.get("properties") or {})
        actual_properties = dict(actual.get("properties") or {})
        property_mismatches = {
            key: {"expected": value, "actual": actual_properties.get(key)}
            for key, value in expected_properties.items()
            if actual_properties.get(key) != value
        }
        if property_mismatches:
            reasons.append("property mismatch")

        status = _comparison_status(reasons)
        entry = {
            "key": expected["key"],
            "status": status,
            "reasons": reasons,
            "expected": expected,
            "actual": actual,
            "propertyMismatches": property_mismatches,
            "fixable": status in {"match", "name_mismatch", "parent_mismatch", "stereotype_mismatch", "property_mismatch"},
            "suggestedActions": _suggested_actions(expected, actual, reasons),
        }
        entries.append(entry)
        if status != "match":
            patch_plan.extend(
                {
                    "stepId": f"rubric:{len(patch_plan) + index}",
                    **action,
                    "status": "preview",
                    "reason": ", ".join(reasons) or "comparison mismatch",
                }
                for index, action in enumerate(entry["suggestedActions"], start=1)
            )

    unexpected = []
    for key, actual in actual_index.items():
        if key in expected_index:
            continue
        entry = {
            "key": key,
            "status": "unexpected",
            "reasons": ["artifact not listed in rubric"],
            "expected": None,
            "actual": actual,
            "fixable": False,
            "suggestedActions": [
                {
                    "action": "review_or_remove_artifact",
                    "targetKey": key,
                    "preview": f"Review whether '{actual.get('name') or key}' belongs in the rubric or should be removed.",
                }
            ],
        }
        unexpected.append(entry)

    entries.extend(unexpected)
    patch_plan.extend(
        {
            "stepId": f"rubric:{len(patch_plan) + index}",
            **action,
            "status": "preview",
            "reason": "unexpected artifact",
        }
        for entry in unexpected
        for index, action in enumerate(entry["suggestedActions"], start=1)
    )

    return {
        "expectedArtifactCount": len(expected_items),
        "currentArtifactCount": len(actual_items),
        "ready": all(entry["status"] == "match" for entry in entries if entry["expected"] is not None)
        and not unexpected,
        "missingArtifactKeys": [entry["key"] for entry in entries if entry["status"] == "missing"],
        "unexpectedArtifactKeys": [entry["key"] for entry in unexpected],
        "entries": entries,
        "patchPlan": patch_plan,
    }


def _workflow_context(
    pack_id: str,
    *,
    recipe_id: str | None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None,
) -> tuple[registry.PackDefinition, dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    pack = registry.get_pack(pack_id)
    expected = _build_expected_artifacts(pack_id, recipe_id=recipe_id)
    workflow_guidance = methodology_service.get_workflow_guidance(
        pack_id=pack_id,
        recipe_id=recipe_id,
        recipe_parameters={},
        completed_artifacts=[_coerce_artifact(item, role="current") for item in (current_artifacts or ())],
    )
    comparison = compare_against_expected_artifact_list(expected, current_artifacts)
    return pack, workflow_guidance, expected, comparison["entries"]


def validate_assignment_package(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
) -> dict[str, Any]:
    """Validate a rubric assignment/package in dry-run mode."""
    pack = registry.get_pack(pack_id)
    expected = (
        [_coerce_artifact(item, role="expected") for item in expected_artifacts]
        if expected_artifacts is not None
        else _build_expected_artifacts(pack_id, recipe_id=recipe_id)
    )
    comparison = compare_against_expected_artifact_list(expected, current_artifacts)
    workflow_guidance = methodology_service.get_workflow_guidance(
        pack_id=pack_id,
        recipe_id=recipe_id,
        recipe_parameters={},
        completed_artifacts=[_coerce_artifact(item, role="current") for item in (current_artifacts or ())],
    )
    ready = comparison["ready"] and bool(workflow_guidance["guidance"].get("ready_to_execute"))
    return {
        "pack": pack.to_dict(),
        "scope": "recipe" if recipe_id is not None else "pack",
        "recipeId": recipe_id,
        "workflowGuidance": workflow_guidance["guidance"],
        "expectedArtifacts": expected,
        "artifactComparison": comparison,
        "ready": ready,
        "recommendedActions": list(workflow_guidance["guidance"].get("recommended_actions", ()))
        + [action["preview"] for entry in comparison["entries"] for action in entry["suggestedActions"]],
        "dryRun": True,
    }


def export_required_diagrams(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    export_format: str = "png",
) -> dict[str, Any]:
    """Build a dry-run export queue for required diagrams."""
    pack = registry.get_pack(pack_id)
    expected = (
        [_coerce_artifact(item, role="expected") for item in expected_artifacts]
        if expected_artifacts is not None
        else _build_expected_artifacts(pack_id, recipe_id=recipe_id)
    )
    comparison = compare_against_expected_artifact_list(expected, current_artifacts)
    export_items: list[dict[str, Any]] = []
    for index, artifact in enumerate(_expected_diagrams(expected), start=1):
        entry = next((item for item in comparison["entries"] if item["key"] == artifact["key"]), None)
        actual = entry["actual"] if entry is not None else None
        export_items.append(
            {
                "stepId": f"export:{index}",
                "artifactKey": artifact["key"],
                "diagramName": artifact.get("name") or artifact["key"],
                "diagramKind": artifact["kind"],
                "recipeId": artifact.get("recipe_id"),
                "status": "ready" if entry is not None and entry["status"] == "match" else "blocked",
                "diagramId": actual.get("element_id") if actual else None,
                "outputFile": f"{_slugify(artifact.get('name') or artifact['key'])}.{export_format.lstrip('.')}",
                "format": export_format.lstrip("."),
                "reasons": list(entry["reasons"]) if entry is not None else ["artifact not found"],
            }
        )

    return {
        "pack": pack.to_dict(),
        "scope": "recipe" if recipe_id is not None else "pack",
        "recipeId": recipe_id,
        "exportFormat": export_format.lstrip("."),
        "expectedDiagrams": _expected_diagrams(expected),
        "artifactComparison": comparison,
        "exportItems": export_items,
        "readyToExport": all(item["status"] == "ready" for item in export_items),
        "blockedDiagramKeys": [item["artifactKey"] for item in export_items if item["status"] != "ready"],
        "dryRun": True,
    }


def assemble_ppt_pdf(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    export_plan: Mapping[str, Any] | None = None,
    title: str | None = None,
    pptx_name: str | None = None,
    pdf_name: str | None = None,
    export_format: str = "png",
) -> dict[str, Any]:
    """Create a dry-run assembly plan for a PowerPoint + PDF submission."""
    pack = registry.get_pack(pack_id)
    export_result = (
        dict(export_plan)
        if export_plan is not None
        else export_required_diagrams(
            pack_id,
            recipe_id=recipe_id,
            current_artifacts=current_artifacts,
            expected_artifacts=expected_artifacts,
            export_format=export_format,
        )
    )
    export_items = list(export_result.get("exportItems") or [])
    deck_title = title or pack.title
    base_name = _slugify(deck_title)
    pptx_file = pptx_name or f"{base_name}.pptx"
    pdf_file = pdf_name or f"{base_name}.pdf"
    slide_items: list[dict[str, Any]] = [
        {
            "slideNumber": 1,
            "kind": "title",
            "title": deck_title,
            "subtitle": f"{pack.domain} rubric workflow",
            "notes": "Cover slide generated from the rubric workflow plan.",
        }
    ]
    for index, item in enumerate(export_items, start=2):
        slide_items.append(
            {
                "slideNumber": index,
                "kind": "diagram",
                "title": item["diagramName"],
                "artifactKey": item["artifactKey"],
                "diagramId": item.get("diagramId"),
                "sourceFile": item["outputFile"],
                "status": item["status"],
                "notes": "; ".join(item.get("reasons") or ()),
            }
        )
    slide_items.append(
        {
            "slideNumber": len(slide_items) + 1,
            "kind": "appendix",
            "title": "Rubric Comparison",
            "notes": "Summary of missing, unexpected, and blocked artifacts.",
        }
    )
    ready = bool(export_result.get("readyToExport"))
    return {
        "pack": pack.to_dict(),
        "scope": "recipe" if recipe_id is not None else "pack",
        "recipeId": recipe_id,
        "presentationPlan": {
            "title": deck_title,
            "pptxFile": pptx_file,
            "pdfFile": pdf_file,
            "slideCount": len(slide_items),
            "slides": slide_items,
            "readyToAssemble": ready,
            "blockedDiagramKeys": list(export_result.get("blockedDiagramKeys") or ()),
        },
        "exportPlan": export_result,
        "dryRun": True,
    }


def _artifact_snapshot_to_dict(artifact: Any) -> dict[str, Any]:
    if hasattr(artifact, "__dataclass_fields__"):
        payload = asdict(artifact)
    elif isinstance(artifact, Mapping):
        payload = dict(artifact)
    else:
        raise TypeError(f"Unsupported artifact snapshot type: {type(artifact)!r}")
    payload["element_id"] = payload.get("element_id") or payload.get("elementId")
    payload["parent_key"] = payload.get("parent_key") or payload.get("parentKey")
    return payload


async def discover_current_artifacts(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    root_package_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    recipe_parameters: Mapping[str, Any] | None = None,
    bridge: Any = default_bridge_client,
) -> list[dict[str, Any]]:
    """Discover live artifacts for a pack or recipe using the methodology service internals."""
    params = dict(recipe_parameters or {})
    pack = registry.get_pack(pack_id)
    runtime_pack = methodology_service._build_runtime_pack(  # type: ignore[attr-defined]
        pack,
        params,
        target_recipe_id=recipe_id,
        strict=False,
    )
    discovered = methodology_service._seed_completed_artifacts(  # type: ignore[attr-defined]
        root_package_id or params.get("root_package_id") or "",
        current_artifacts or (),
    )
    if not root_package_id and params.get("root_package_id"):
        root_package_id = str(params["root_package_id"])

    for runtime_recipe in runtime_pack.recipes:
        references = methodology_service._reference_artifacts(  # type: ignore[attr-defined]
            runtime_recipe.recipe_id,
            params,
        )
        existing_keys = {artifact.key for artifact in discovered}
        discovered = list(discovered) + [
            reference for reference in references if reference.key not in existing_keys
        ]
        discovered = await methodology_service._discover_live_artifacts(  # type: ignore[attr-defined]
            runtime_recipe,
            discovered,
            recipe_parameters=params,
            root_package_id=root_package_id,
            bridge=bridge,
        )

    return [_artifact_snapshot_to_dict(artifact) for artifact in discovered]


def _decode_exported_image(payload: Mapping[str, Any]) -> bytes:
    image = payload.get("image")
    if not isinstance(image, str) or not image:
        raise ValueError("bridge did not return a base64-encoded diagram image")
    return base64.b64decode(image)


def _write_pdf(image_paths: Sequence[Path], output_path: Path) -> None:
    if not image_paths:
        raise ValueError("cannot assemble PDF without exported images")

    images: list[Image.Image] = []
    try:
        for image_path in image_paths:
            image = Image.open(image_path)
            images.append(image.convert("RGB"))
        first, rest = images[0], images[1:]
        first.save(output_path, save_all=True, append_images=rest)
    finally:
        for image in images:
            image.close()


def _write_pptx(
    image_paths: Sequence[Path],
    output_path: Path,
    *,
    title: str,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:  # pragma: no cover - depends on optional dependency
        raise RuntimeError(
            "python-pptx is required to assemble PPTX output. Install the mcp-server package dependencies."
        ) from exc

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    cover = presentation.slides.add_slide(title_layout)
    cover.shapes.title.text = title
    subtitle = cover.placeholders[1] if len(cover.placeholders) > 1 else None
    if subtitle is not None:
        subtitle.text = "Generated by Cameo MCP rubric workflow"

    blank_layout = presentation.slide_layouts[6]
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    for image_path in image_paths:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=slide_width,
            height=slide_height,
        )
        textbox = slide.shapes.add_textbox(
            Inches(0.2),
            Inches(0.1),
            Inches(9.4),
            Inches(0.4),
        )
        textbox.text_frame.text = image_path.stem.replace("-", " ").title()

    presentation.save(output_path)


async def validate_assignment_package_live(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    root_package_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    bridge: Any = default_bridge_client,
) -> dict[str, Any]:
    discovered = (
        list(current_artifacts)
        if current_artifacts is not None
        else await discover_current_artifacts(
            pack_id,
            recipe_id=recipe_id,
            root_package_id=root_package_id,
            bridge=bridge,
        )
    )
    result = validate_assignment_package(
        pack_id,
        recipe_id=recipe_id,
        current_artifacts=discovered,
        expected_artifacts=expected_artifacts,
    )
    result["currentArtifacts"] = discovered
    result["dryRun"] = True
    return result


async def export_required_diagrams_live(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    root_package_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    export_format: str = "png",
    output_dir: str | None = None,
    bridge: Any = default_bridge_client,
) -> dict[str, Any]:
    discovered = (
        list(current_artifacts)
        if current_artifacts is not None
        else await discover_current_artifacts(
            pack_id,
            recipe_id=recipe_id,
            root_package_id=root_package_id,
            bridge=bridge,
        )
    )
    plan = export_required_diagrams(
        pack_id,
        recipe_id=recipe_id,
        current_artifacts=discovered,
        expected_artifacts=expected_artifacts,
        export_format=export_format,
    )
    plan["currentArtifacts"] = discovered
    if output_dir is None:
        return plan

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    exports: list[dict[str, Any]] = []
    for item in plan["exportItems"]:
        if item["status"] != "ready" or not item.get("diagramId"):
            exports.append(
                {
                    "artifactKey": item["artifactKey"],
                    "status": "blocked",
                    "reasons": item.get("reasons") or ["diagram not ready"],
                    "outputFile": str(out_dir / item["outputFile"]),
                }
            )
            continue
        payload = await bridge.get_diagram_image(
            item["diagramId"],
            format=export_format,
        )
        image_bytes = _decode_exported_image(payload)
        target_path = out_dir / item["outputFile"]
        target_path.write_bytes(image_bytes)
        exports.append(
            {
                "artifactKey": item["artifactKey"],
                "diagramId": item["diagramId"],
                "status": "exported",
                "outputFile": str(target_path),
                "format": export_format.lstrip("."),
                "width": payload.get("width"),
                "height": payload.get("height"),
                "byteCount": len(image_bytes),
            }
        )

    plan["dryRun"] = False
    plan["outputDir"] = str(out_dir)
    plan["exports"] = exports
    plan["exportedCount"] = sum(1 for item in exports if item["status"] == "exported")
    plan["blockedCount"] = sum(1 for item in exports if item["status"] != "exported")
    return plan


async def assemble_ppt_pdf_live(
    pack_id: str,
    *,
    recipe_id: str | None = None,
    root_package_id: str | None = None,
    current_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    expected_artifacts: Sequence[Mapping[str, Any] | Any] | None = None,
    output_dir: str | None = None,
    title: str | None = None,
    pptx_name: str | None = None,
    pdf_name: str | None = None,
    export_format: str = "png",
    bridge: Any = default_bridge_client,
) -> dict[str, Any]:
    export_result = await export_required_diagrams_live(
        pack_id,
        recipe_id=recipe_id,
        root_package_id=root_package_id,
        current_artifacts=current_artifacts,
        expected_artifacts=expected_artifacts,
        export_format=export_format,
        output_dir=output_dir,
        bridge=bridge,
    )
    plan = assemble_ppt_pdf(
        pack_id,
        recipe_id=recipe_id,
        current_artifacts=export_result.get("currentArtifacts"),
        expected_artifacts=expected_artifacts,
        export_plan=export_result,
        title=title,
        pptx_name=pptx_name,
        pdf_name=pdf_name,
        export_format=export_format,
    )
    if output_dir is None:
        return plan

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    exported_paths = [
        Path(item["outputFile"])
        for item in export_result.get("exports") or ()
        if item.get("status") == "exported"
    ]
    pdf_path = out_dir / plan["presentationPlan"]["pdfFile"]
    pptx_path = out_dir / plan["presentationPlan"]["pptxFile"]
    created: list[dict[str, Any]] = []
    if exported_paths:
        _write_pdf(exported_paths, pdf_path)
        created.append({"kind": "pdf", "path": str(pdf_path)})
        _write_pptx(exported_paths, pptx_path, title=plan["presentationPlan"]["title"])
        created.append({"kind": "pptx", "path": str(pptx_path)})

    plan["dryRun"] = False
    plan["outputDir"] = str(out_dir)
    plan["createdFiles"] = created
    plan["ready"] = bool(exported_paths)
    return plan


__all__ = [
    "assemble_ppt_pdf",
    "assemble_ppt_pdf_live",
    "compare_against_expected_artifact_list",
    "discover_current_artifacts",
    "export_required_diagrams",
    "export_required_diagrams_live",
    "validate_assignment_package",
    "validate_assignment_package_live",
]
